"""
Text Extractor Module

Handles text extraction from various document formats with support for:
- Multiple file formats
- OCR for images and scanned documents
- Structured data extraction from spreadsheets and presentations
- Malicious content detection
"""

import os
import zipfile
import requests
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import pandas as pd
from docx import Document
from pptx import Presentation
import io
import mimetypes
import hashlib
from typing import Dict, Any, Tuple, Optional
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextExtractor:
    """Main text extraction class with support for multiple formats"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def is_malicious_content(self, content: str) -> bool:
        """Check if content contains potentially malicious patterns"""
        malicious_patterns = [
            'javascript:', 'vbscript:', 'onload=', 'onerror=',
            'eval(', 'exec(', 'system(', 'shell_exec('
        ]
        content_lower = content.lower()
        return any(pattern in content_lower for pattern in malicious_patterns)
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF with OCR fallback"""
        try:
            doc = fitz.open(file_path)
            text = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                
                if not page_text.strip():
                    # Try OCR if no text found
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    page_text = pytesseract.image_to_string(img)
                
                text += page_text + "\n"
            
            doc.close()
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return ""
    
    def extract_text_from_word(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from Word document {file_path}: {e}")
            return ""
    
    def extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel spreadsheet"""
        try:
            df = pd.read_excel(file_path, sheet_name=None)
            text = ""
            for sheet_name, sheet_df in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_df.to_string(index=False) + "\n\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from Excel file {file_path}: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"Slide {slide_num + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint file {file_path}: {e}")
            return ""
    
    def extract_text_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from image {file_path}: {e}")
            return ""
    
    def extract_text_from_html(self, file_path: str) -> str:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = ' '.join(chunk for chunk in chunks if chunk)
                return text
        except Exception as e:
            logger.error(f"Error extracting text from HTML file {file_path}: {e}")
            return ""
    
    def extract_text_from_zip(self, file_path: str) -> str:
        """Extract text from ZIP archive"""
        try:
            text = ""
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                for file_info in zip_ref.filelist:
                    if not file_info.is_dir():
                        file_name = file_info.filename
                        if any(ext in file_name.lower() for ext in ['.txt', '.md', '.csv']):
                            with zip_ref.open(file_name) as f:
                                text += f"File: {file_name}\n"
                                text += f.read().decode('utf-8', errors='ignore') + "\n\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from ZIP file {file_path}: {e}")
            return ""
    
    def extract_text_from_binary(self, file_path: str) -> str:
        """Extract text from binary file (basic hex dump)"""
        try:
            with open(file_path, 'rb') as f:
                data = f.read(1024)  # Read first 1KB
                hex_dump = data.hex()
                return f"Binary file hex dump (first 1KB): {hex_dump[:200]}..."
        except Exception as e:
            logger.error(f"Error reading binary file {file_path}: {e}")
            return ""
    
    def detect_file_type_from_content(self, file_path: str) -> str:
        """Detect file type based on content and extension"""
        try:
            # Check file extension first
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()
            
            # Known extensions
            if ext in ['.pdf']:
                return 'pdf'
            elif ext in ['.doc', '.docx']:
                return 'word'
            elif ext in ['.xls', '.xlsx']:
                return 'excel'
            elif ext in ['.ppt', '.pptx']:
                return 'pptx'
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                return 'image'
            elif ext in ['.html', '.htm']:
                return 'html'
            elif ext in ['.zip']:
                return 'zip'
            else:
                # Try to detect by content
                with open(file_path, 'rb') as f:
                    header = f.read(8)
                    
                if header.startswith(b'%PDF'):
                    return 'pdf'
                elif header.startswith(b'PK'):
                    return 'zip'
                elif header.startswith(b'\xff\xd8\xff'):
                    return 'image'
                else:
                    return 'unknown'
                    
        except Exception as e:
            logger.error(f"Error detecting file type for {file_path}: {e}")
            return 'unknown'
    
    def extract_from_file(self, file_path: str) -> Tuple[str, str]:
        """Main extraction method that automatically detects format and extracts text"""
        try:
            file_type = self.detect_file_type_from_content(file_path)
            logger.info(f"Detected file type: {file_type} for {file_path}")
            
            if file_type == 'pdf':
                text = self.extract_text_from_pdf(file_path)
            elif file_type == 'word':
                text = self.extract_text_from_word(file_path)
            elif file_type == 'excel':
                text = self.extract_text_from_excel(file_path)
            elif file_type == 'pptx':
                text = self.extract_text_from_pptx(file_path)
            elif file_type == 'image':
                text = self.extract_text_from_image(file_path)
            elif file_type == 'html':
                text = self.extract_text_from_html(file_path)
            elif file_type == 'zip':
                text = self.extract_text_from_zip(file_path)
            else:
                text = self.extract_text_from_binary(file_path)
            
            # Check for malicious content
            if self.is_malicious_content(text):
                logger.warning(f"Potentially malicious content detected in {file_path}")
                text = "[MALICIOUS CONTENT DETECTED - EXTRACTION BLOCKED]"
            
            return text, file_type
            
        except Exception as e:
            logger.error(f"Error extracting from file {file_path}: {e}")
            return "", "error"

# Standalone functions for backward compatibility
def extract_text_from_pdf(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_pdf(file_path)

def extract_text_from_word(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_word(file_path)

def extract_text_from_excel(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_excel(file_path)

def extract_text_from_pptx(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_pptx(file_path)

def extract_text_from_image(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_image(file_path)

def extract_text_from_html(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_html(file_path)

def extract_text_from_zip(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_zip(file_path)

def extract_text_from_binary(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.extract_text_from_binary(file_path)

def detect_file_type_from_content(file_path: str) -> str:
    extractor = TextExtractor()
    return extractor.detect_file_type_from_content(file_path)
